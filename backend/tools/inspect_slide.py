"""Örnek .pptm dosyalarını inceleme aracı (Faz 0 kalibrasyon).

Kullanım:
    python inspect_slide.py "<dosya yolu>"
"""
import sys
from pptx import Presentation
from pptx.util import Emu


def describe_shape(shape, depth=0):
    indent = "  " * depth
    info = f"{indent}- id={shape.shape_id} name={shape.name!r} type={shape.shape_type}"
    try:
        info += f" pos=({Emu(shape.left).inches:.2f}in,{Emu(shape.top).inches:.2f}in)"
        info += f" size=({Emu(shape.width).inches:.2f}in x {Emu(shape.height).inches:.2f}in)"
    except Exception:
        pass
    print(info)
    if shape.has_text_frame:
        text = shape.text_frame.text.strip().replace("\n", " | ")
        if text:
            print(f"{indent}  text: {text[:120]}")
    if shape.shape_type == 6:  # GROUP
        for sub in shape.shapes:
            describe_shape(sub, depth + 1)


def main(path):
    prs = Presentation(path)
    print(f"Dosya: {path}")
    print(f"Slayt boyutu: {prs.slide_width} x {prs.slide_height} EMU")
    print(f"Slayt sayısı: {len(prs.slides.__iter__.__self__._sldIdLst) if False else len(prs.slides._sldIdLst)}")
    for idx, slide in enumerate(prs.slides, start=1):
        print(f"\n=== Slayt {idx} (layout={slide.slide_layout.name}) ===")
        for shape in slide.shapes:
            describe_shape(shape)


if __name__ == "__main__":
    main(sys.argv[1])
