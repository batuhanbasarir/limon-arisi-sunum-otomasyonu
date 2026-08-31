"""Marka + sıralı içerik listesi verildiğinde tam bir .pptx üretir.

Kapak → her içerik öğesi için içerik slaydı → kapanış sırasıyla,
templates/_shared altındaki paylaşılan skeleton/fragment'leri ve
templates/<brand>/brand.json'daki marka adını kullanır.

Kapak/kapanış/içerik-dekor shape ağaçları üç örnek dosyada da (Erişun,
Miluni, Nudo) MD5 düzeyinde birebir aynı olduğu Faz 0'da doğrulandığı için
tek bir paylaşılan şablon kullanılıyor; markalar arası tek fark kapak
slaydındaki marka adı metni.
"""
import json
from dataclasses import dataclass, field
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Emu, Pt, Inches

from . import slide_cloner
from .pptx_repair import fix_pptx_package
from ..paths import get_project_root

PROJECT_ROOT = get_project_root()
TEMPLATES_DIR = PROJECT_ROOT / "templates"
SHARED_DIR = TEMPLATES_DIR / "_shared"
BADGES_DIR = SHARED_DIR / "badges"
LOGO_PATH = BADGES_DIR / "limon_arisi_logo.png"

CONTENT_LAYOUT_INDEX = 1  # "Başlık ve İçerik" — slideLayout2.xml

_layout_config = json.loads((SHARED_DIR / "layout.json").read_text(encoding="utf-8"))


@dataclass
class ContentItem:
    caption: str
    image_paths: list[str] = field(default_factory=list)
    video_path: str | None = None
    poster_path: str | None = None


def _load_brand(brand_id: str) -> dict:
    path = TEMPLATES_DIR / brand_id / "brand.json"
    if not path.exists():
        raise ValueError(f"Bilinmeyen marka: {brand_id}")
    return json.loads(path.read_text(encoding="utf-8"))


def _box(name: str):
    cfg = _layout_config[name]
    return (
        Inches(cfg["left_in"]),
        Inches(cfg["top_in"]),
        Inches(cfg["width_in"]),
        Inches(cfg["height_in"]),
    )


def _fit_within_box(img_w: int, img_h: int, box: tuple[int, int, int, int]):
    """Görseli, oranını bozmadan (contain-fit) kutunun içine ortalayarak
    sığdırır. `çekiştirme` sorununun kök çözümü — asla width/height'ı
    bağımsız olarak kutuya zorlamayız."""
    left, top, width, height = box
    box_ratio = width / height
    img_ratio = img_w / img_h
    if img_ratio > box_ratio:
        new_width = width
        new_height = round(width / img_ratio)
    else:
        new_height = height
        new_width = round(height * img_ratio)
    offset_x = left + (width - new_width) // 2
    offset_y = top + (height - new_height) // 2
    return offset_x, offset_y, new_width, new_height


def _add_picture_contained(slide, image_path, box_name: str):
    with Image.open(image_path) as img:
        img_w, img_h = img.size
    left, top, width, height = _fit_within_box(img_w, img_h, _box(box_name))
    slide.shapes.add_picture(str(image_path), left, top, width, height)


def _add_logo(slide, box_name: str):
    _add_picture_contained(slide, LOGO_PATH, box_name)


def _add_caption(slide, text: str, box_name: str):
    left, top, width, height = _box(box_name)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Tarayıcılar <textarea> içeriğini form gönderirken \n'i \r\n'e çevirir
    # (HTML spec). Normalize etmezsek her satırın sonunda kalan \r karakteri
    # PowerPoint XML'inde literal "_x000D_" olarak görünür.
    lines = text.replace("\r\n", "\n").replace("\r", "\n").split("\n")
    tf.paragraphs[0].text = lines[0]
    tf.paragraphs[0].runs[0].font.size = Pt(18)
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        if p.runs:
            p.runs[0].font.size = Pt(18)


def _add_content_slide(prs, layout, item: ContentItem):
    slide = slide_cloner.new_slide(prs, layout)
    decor = slide_cloner.load_fragment(SHARED_DIR / "content_decor.xml.fragment")
    slide_cloner.apply_fragment(slide, decor, BADGES_DIR)
    _add_logo(slide, "logo_small")

    if item.video_path:
        # Poster kare, videonun gerçek kare boyutlarıyla birebir aynı — oranı
        # ondan okuyup videoyu da aynı kutuya sığdırıyoruz.
        with Image.open(item.poster_path) as poster:
            poster_w, poster_h = poster.size
        left, top, width, height = _fit_within_box(poster_w, poster_h, _box("media_video"))
        slide.shapes.add_movie(
            item.video_path, left, top, width, height,
            poster_frame_image=item.poster_path,
            mime_type="video/mp4",
        )
        _add_caption(slide, item.caption, "caption")
    elif len(item.image_paths) == 1:
        _add_picture_contained(slide, item.image_paths[0], "media_1img")
        _add_caption(slide, item.caption, "caption")
    elif len(item.image_paths) == 2:
        for path, box_name in zip(item.image_paths, ("media_2img_left", "media_2img_right")):
            _add_picture_contained(slide, path, box_name)
        _add_caption(slide, item.caption, "caption")
    else:
        raise ValueError("ContentItem için 1-2 image_paths veya video_path gerekli")

    return slide


def build_deck(brand_id: str, items: list[ContentItem]) -> Presentation:
    brand = _load_brand(brand_id)

    prs = Presentation(SHARED_DIR / "skeleton.pptx")
    layout = prs.slide_masters[0].slide_layouts[CONTENT_LAYOUT_INDEX]

    cover_slide = slide_cloner.new_slide(prs, layout)
    cover_frag = slide_cloner.load_fragment(SHARED_DIR / "cover.xml.fragment")
    slide_cloner.apply_fragment(cover_slide, cover_frag, BADGES_DIR, brand_name=brand["display_name"])
    _add_logo(cover_slide, "logo_big")

    for item in items:
        _add_content_slide(prs, layout, item)

    closing_slide = slide_cloner.new_slide(prs, layout)
    closing_frag = slide_cloner.load_fragment(SHARED_DIR / "closing.xml.fragment")
    slide_cloner.apply_fragment(closing_slide, closing_frag, BADGES_DIR)
    _add_logo(closing_slide, "logo_small")

    return prs


def save_deck(prs: Presentation, path: Path):
    """prs.save() + docProps/app.xml ve Content_Types tutarlılık düzeltmesi
    (bkz. pptx_repair.py) — bu adım atlanırsa PowerPoint dosyayı bozuk
    sanıp onarım isteyebilir."""
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)
    fix_pptx_package(path, slide_count=len(prs.slides))
