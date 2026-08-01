"""Dekoratif slayt fragment'lerini (kapak/kapanış/içerik-dekor) yeni
slaytlara klonlayan çekirdek motor.

Fragment dosyaları templates/_shared/*.xml.fragment altında saklanır ve
görsel referanslarını {{IMG:<key>}}, marka adını {{BRAND_NAME}} token'ı
olarak tutar (bkz. backend/tools/extract_brand_profile.py). Bu modül,
token'ları çalışma zamanında gerçek ilişki id'leri / metinlerle değiştirip
shape XML'ini yeni slaydın spTree'sine ekler.
"""
import re
from copy import deepcopy
from pathlib import Path

from lxml import etree
from pptx.oxml.ns import qn

IMG_TOKEN_RE = re.compile(r"\{\{IMG:([a-zA-Z0-9_]+)\}\}")
BRAND_TOKEN = "{{BRAND_NAME}}"

_R_EMBED = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
_R_LINK = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}link"


def load_fragment(path: Path) -> etree._Element:
    return etree.parse(str(path)).getroot()


def new_slide(prs, layout):
    """Layout'tan otomatik miras alınan placeholder shape'leri temizlenmiş
    boş bir slayt oluşturur (fragment'ler kendi shape'lerini getirir)."""
    slide = prs.slides.add_slide(layout)
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            shape._element.getparent().remove(shape._element)
    return slide


def _resolve_image_path(image_dir: Path, key: str) -> Path:
    matches = list(image_dir.glob(f"{key}.*"))
    if not matches:
        raise FileNotFoundError(f"Görsel bulunamadı: {image_dir}/{key}.*")
    return matches[0]


def apply_fragment(slide, fragment_root, image_dir: Path, brand_name: str | None = None):
    """fragment_root içindeki her shape'i deepcopy'leyip slayda ekler;
    {{IMG:key}} ve {{BRAND_NAME}} token'larını çalışma zamanı değerleriyle
    değiştirir."""
    spTree = slide.shapes._spTree
    image_cache: dict[str, str] = {}

    for shape_el in fragment_root:
        shape_copy = deepcopy(shape_el)

        for el in shape_copy.iter():
            for attr in (_R_EMBED, _R_LINK):
                val = el.get(attr)
                if not val:
                    continue
                m = IMG_TOKEN_RE.fullmatch(val)
                if not m:
                    continue
                key = m.group(1)
                if key not in image_cache:
                    img_path = _resolve_image_path(image_dir, key)
                    _part, rId = slide.part.get_or_add_image_part(str(img_path))
                    image_cache[key] = rId
                el.set(attr, image_cache[key])

        if brand_name is not None:
            for t in shape_copy.iter(qn("a:t")):
                if t.text and BRAND_TOKEN in t.text:
                    t.text = t.text.replace(BRAND_TOKEN, brand_name)

        spTree.append(shape_copy)
